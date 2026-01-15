import os
import fitz  # PyMuPDF
import docx
from pptx import Presentation
import pandas as pd
import chromadb
from chromadb.config import Settings
from sentence_transformers import SentenceTransformer
import uuid
import datetime

class RAGEngine:
    def __init__(self, persist_directory="./vector_db"):
        print("[RAG] Initializing RAG Engine...")
        
        # Initialize Vector DB (Chroma)
        self.chroma_client = chromadb.PersistentClient(path=persist_directory)
        self.collection = self.chroma_client.get_or_create_collection(name="learnpilot_docs")
        
        # Initialize Embedding Model (Local & Free)
        print("[RAG] Loading Embedding Model (all-MiniLM-L6-v2)...")
        self.embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
        print("[RAG] Initialization Complete.")

    def extract_text(self, file_path, original_name):
        """Extracts text from PDF, DOCX, PPTX, TXT, CSV, XLS."""
        ext = os.path.splitext(original_name)[1].lower()
        text = ""

        try:
            if ext == '.pdf':
                with fitz.open(file_path) as doc:
                    for page in doc:
                        text += page.get_text() + "\n"
            
            elif ext == '.docx':
                doc = docx.Document(file_path)
                text = "\n".join([para.text for para in doc.paragraphs])
            
            elif ext == '.pptx':
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            
            elif ext == '.txt':
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()
            
            elif ext in ['.csv', '.xlsx', '.xls']:
                # For tabular data, we convert rows to text strings
                if ext == '.csv':
                    df = pd.read_csv(file_path)
                else:
                    df = pd.read_excel(file_path)
                
                # Convert dataframe to string representation
                text = df.to_string(index=False)
            
            elif ext in ['.png', '.jpg', '.jpeg']:
                # Lazy load EasyOCR only when needed to save startup time
                import easyocr
                reader = easyocr.Reader(['en']) 
                results = reader.readtext(file_path)
                # Combine detected text
                text = " ".join([res[1] for res in results])
            
            else:
                raise ValueError(f"Unsupported file extension: {ext}")
            
            return text.strip()

        except Exception as e:
            print(f"[RAG] Extraction Error for {original_name}: {e}")
            raise e

    def chunk_text(self, text, chunk_size=1000, overlap=200):
        """Splits text into chunks with overlap, respecting word boundaries."""
        if not text:
            return []
        
        chunks = []
        start = 0
        text_len = len(text)

        while start < text_len:
            end = start + chunk_size
            
            # If we are not at the end of text, try to find a space to break at
            if end < text_len:
                # Look for the last space within the chunk to avoid cutting words
                last_space = text.rfind(' ', start, end)
                if last_space != -1:
                    end = last_space
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            # Should move forward, but ensure we don't get stuck if no space found (fallback to hard cut)
            if end == start: 
                end = start + chunk_size
            
            start = end - overlap
            if start < 0: start = 0 # Safety
        
        return chunks

    def ingest_file(self, file_path, file_id, user_id, original_name):
        """Main pipeline: Extract -> Chunk -> Embed -> Store."""
        print(f"[RAG] Ingesting: {original_name} ({file_id})")
        
        # 1. Extract
        raw_text = self.extract_text(file_path, original_name)
        if not raw_text:
            return {"status": "failed", "reason": "No text extracted"}

        # 2. Determine Dynamic Chunk Size
        ext = os.path.splitext(original_name)[1].lower()
        if ext in ['.csv', '.xlsx', '.xls']:
            chunk_size = 2000 # Larger chunks for spreadsheets to capture more rows
            overlap = 400
        elif ext in ['.pdf', '.docx', '.pptx', '.png', '.jpg', '.jpeg']:
            chunk_size = 1500 # Standard page-like content
            overlap = 300
        else:
            chunk_size = 1200 # Default for txt, etc.
            overlap = 200

        # 3. Chunk
        text_chunks = self.chunk_text(raw_text, chunk_size=chunk_size, overlap=overlap)
        print(f"[RAG] Generated {len(text_chunks)} chunks (Size: {chunk_size}).")

        # 4. Embed
        embeddings = self.embedding_model.encode(text_chunks).tolist()

        # 5. Store in Chroma
        ids = [f"{file_id}_{i}" for i in range(len(text_chunks))]
        metadatas = [{
            "userId": user_id,
            "fileId": file_id,
            "source": original_name,
            "chunkIndex": i,
            "timestamp": str(datetime.datetime.now())
        } for i in range(len(text_chunks))]

        self.collection.add(
            documents=text_chunks,
            embeddings=embeddings,
            metadatas=metadatas,
            ids=ids
        )

        return {
            "status": "success",
            "chunks_count": len(text_chunks),
            "file_id": file_id
        }

    def query(self, user_id, query_text, k=5):
        """Retrieves relevant context for a query."""
        print(f"[RAG] Querying for User {user_id}: {query_text}")
        
        # 1. Embed Query
        query_embedding = self.embedding_model.encode([query_text]).tolist()

        # 2. Search Chroma (Filter by userId)
        results = self.collection.query(
            query_embeddings=query_embedding,
            n_results=k,
            where={"userId": user_id} # Strict data isolation
        )

        context_chunks = results['documents'][0]
        sources = [m['source'] for m in results['metadatas'][0]]

        # 3. Format Response
        if not context_chunks:
            return {
                "answer": "I couldn't find any relevant information in your uploaded materials.",
                "sources": []
            }

        # (Optional) Here you would pass 'context_chunks' to an LLM to generate the final answer.
        # Since we are sticking to "Free/Local" without a heavy heavy local LLM requirement in the prompt 
        # (user asked for RAG pipeline, but embeddings is one thing, running Llama-7B locally is another 
        # and might be heavy. The user *did* mention 'Chatbot pipeline: RAG', but only specified 'sentence-transformers' 
        # as the model. I will return the GROUNDED CONTEXT so the Node/Frontend can simplistically show it 
        # or the user can assume an LLM generation step happens here if they had one.)
        
        # Construct a simple grounded response for now:
        combined_context = "\n\n".join(context_chunks)
        
        return {
            "answer": f"Based on your uploaded documents, here is the relevant information:\n\n{combined_context}",
            "sources": list(set(sources)), # Unique sources
            "raw_context": context_chunks
        }
